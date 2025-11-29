
# main.py
"""
Population Dashboard (based on your layout) with:
- same sidebar & logic you provided
- export PDF / PPTX (safe on Windows)
- "Phân tích chuyên sâu (AI)" tab supporting OpenAI GPT and Google Gemini (if SDK available)
- loading spinners and UI report displayed on page
"""

import streamlit as st
import pandas as pd
import requests
import plotly.express as px
from datetime import datetime
import tempfile
import os
import time
from io import BytesIO

# PDF
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER
from reportlab.lib import colors
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont

# PPTX
from pptx import Presentation
from pptx.util import Inches

# Try optional AI SDKs
try:
    import openai
    OPENAI_AVAILABLE = True
except Exception:
    OPENAI_AVAILABLE = False

try:
    import google.generativeai as genai
    GEMINI_AVAILABLE = True
except Exception:
    GEMINI_AVAILABLE = False

# Register Chinese/Vietnamese-capable CID font to reduce encoding problems in PDF
# (ReportLab has limited TTF support without extra setup; UnicodeCIDFont 'MSung-Light' is a common CID)
try:
    pdfmetrics.registerFont(UnicodeCIDFont("MSung-Light"))
except Exception:
    # ignore if not available; ReportLab will fallback
    pass

st.set_page_config(page_title="Population Dashboard", layout="wide")
st.title("📊 Phân tích dân số theo tỉ lệ sinh - tử (Extended)")
st.write("Chọn 1 hoặc 2 quốc gia để so sánh; chọn khoảng thời gian; xuất PDF/PPTX từ báo cáo.")

WB_BASE = "http://api.worldbank.org/v2"

# ---------------- Data functions (unchanged logic)
def get_country_list_worldbank():
    url = f"{WB_BASE}/country?format=json&per_page=500"
    res = requests.get(url, timeout=10)
    data = res.json()
    countries = []
    if isinstance(data, list) and len(data) >= 2:
        for c in data[1]:
            countries.append({"id": c.get("id"), "name": c.get("name"), "iso2Code": c.get("iso2Code")})
    return countries


def _fetch_indicator(country_id, indicator, start_year, end_year):
    url = f"{WB_BASE}/country/{country_id}/indicator/{indicator}?date={start_year}:{end_year}&format=json&per_page=1000"
    res = requests.get(url, timeout=15)
    data = res.json()
    series = {}
    if isinstance(data, list) and len(data) >= 2:
        for item in data[1]:
            try:
                year = int(item.get("date"))
                value = item.get("value")
                if value is not None:
                    series[year] = float(value)
            except Exception:
                continue
    return series


def get_series_for_country(country_id, start_year, end_year):
    """Return DataFrame with columns: year, birth_rate, death_rate, population"""
    birth_ind = "SP.DYN.CBRT.IN"
    death_ind = "SP.DYN.CDRT.IN"
    pop_ind = "SP.POP.TOTL"
    b = _fetch_indicator(country_id, birth_ind, start_year, end_year)
    d = _fetch_indicator(country_id, death_ind, start_year, end_year)
    p = _fetch_indicator(country_id, pop_ind, start_year, end_year)
    years = sorted(set(list(b.keys()) + list(d.keys()) + list(p.keys())))
    rows = []
    for y in years:
        rows.append({"year": y, "birth_rate": b.get(y), "death_rate": d.get(y), "population": p.get(y)})
    if not rows:
        return pd.DataFrame()
    df = pd.DataFrame(rows).sort_values("year")
    return df

# ---------------- Sidebar (keep your UI)
st.sidebar.header("Tùy chọn")
countries = get_country_list_worldbank()
df_countries = pd.DataFrame(countries)

sel_names = st.sidebar.multiselect(
    "Chọn quốc gia (tối đa 2):",
    df_countries['name'].tolist(),
    default=["Viet Nam"] if "Viet Nam" in df_countries['name'].tolist() else []
)

current_year = datetime.now().year
start_year = st.sidebar.slider("Năm bắt đầu:", min_value=1960, max_value=current_year - 1, value=2015)
end_year = st.sidebar.slider("Năm kết thúc:", min_value=start_year + 1, max_value=current_year, value=current_year)

# AI quick options in sidebar
st.sidebar.markdown("---")
st.sidebar.markdown("**Phân tích chuyên sâu (AI)**")
model_default = "openai"  # default
model_choice_sidebar = st.sidebar.selectbox("Model mặc định khi mở tab AI:", ["openai", "gemini"])
st.sidebar.caption("Bạn có thể đổi model trong tab 'Phân tích chuyên sâu'.")

openai_key_sidebar = None
if model_choice_sidebar == "openai":
    if OPENAI_AVAILABLE:
        openai_key_sidebar = st.sidebar.text_input("OpenAI API key (tùy chọn)", type="password")
    else:
        st.sidebar.warning("openai SDK chưa cài. Cài `openai` nếu muốn dùng OpenAI.")
else:
    if GEMINI_AVAILABLE:
        gemini_key_sidebar = st.sidebar.text_input("Gemini API key (tùy chọn)", type="password")
    else:
        st.sidebar.info("Nếu muốn dùng Gemini, cài `google-generativeai` và nhập key trong tab AI.")

# ---------------- Load data with spinner
if not sel_names:
    st.warning("Vui lòng chọn ít nhất 1 quốc gia ở thanh bên để bắt đầu.")
    st.stop()

with st.spinner("⏳ Đang tải dữ liệu từ World Bank..."):
    time.sleep(0.8)
    country_dfs = {}
    for name in sel_names:
        cid = df_countries[df_countries['name'] == name]['id'].values[0]
        df = get_series_for_country(cid, start_year, end_year)
        country_dfs[name] = df

any_data = any([not df.empty for df in country_dfs.values()])
if not any_data:
    st.warning("Không tìm thấy dữ liệu cho quốc gia/khoảng thời gian đã chọn.")
    st.stop()

# ---------------- Main content (keep charts & tables)
st.subheader("📋 Dữ liệu (bảng)")
for name, df in country_dfs.items():
    if df.empty:
        st.write(f"**{name}**: Không có dữ liệu.")
    else:
        st.write(f"**{name}**")
        st.dataframe(df)

st.subheader("📈 So sánh tỉ lệ sinh & tử")
plot_df = pd.DataFrame()
for name, df in country_dfs.items():
    if not df.empty:
        tmp = df.copy()
        tmp['country'] = name
        plot_df = pd.concat([plot_df, tmp], ignore_index=True)

if not plot_df.empty:
    fig = px.line(plot_df, x='year', y=['birth_rate', 'death_rate'], color='country',
                  labels={'value': 'Tỉ lệ (per 1,000)'})
    st.plotly_chart(fig, use_container_width=True)

st.subheader("📊 Tương quan: Population vs Birth Rate")
scatter_df = pd.DataFrame()
for name, df in country_dfs.items():
    if not df.empty:
        tmp = df[['year', 'birth_rate', 'population']].copy()
        tmp['country'] = name
        scatter_df = pd.concat([scatter_df, tmp], ignore_index=True)

if not scatter_df.empty:
    fig2 = px.scatter(scatter_df, x='population', y='birth_rate', color='country', size='population',
                      hover_data=['year'])
    fig2.update_layout(xaxis_title='Population', yaxis_title='Birth rate (per 1,000)')
    st.plotly_chart(fig2, use_container_width=True)

# ---------------- Auto report table display
st.subheader("📑 Phân tích tự động (tóm tắt)")

report_rows = []
for name, df in country_dfs.items():
    if df.empty:
        continue
    avg_b = df['birth_rate'].mean()
    avg_d = df['death_rate'].mean()
    trend = '📉 Giảm' if df['birth_rate'].iloc[-1] < df['birth_rate'].iloc[0] else '📈 Tăng'
    report_rows.append({
        "Quốc gia": name,
        "Tỉ lệ sinh TB (‰)": f"{avg_b:.2f}",
        "Tỉ lệ tử TB (‰)": f"{avg_d:.2f}",
        "Xu hướng tỉ lệ sinh": trend
    })

if report_rows:
    report_df = pd.DataFrame(report_rows)
    st.table(report_df)
else:
    st.info("Không đủ dữ liệu để tạo báo cáo.")

# Show detailed textual report on page (clean report layout)
st.markdown("### 🧾 Báo cáo chi tiết (hiển thị trên giao diện)")
for r in report_rows:
    st.markdown(f"#### 🌍 {r['Quốc gia']}")
    st.markdown(f"- **Tỉ lệ sinh trung bình:** {r['Tỉ lệ sinh TB (‰)']} ‰")
    st.markdown(f"- **Tỉ lệ tử trung bình:** {r['Tỉ lệ tử TB (‰)']} ‰")
    st.markdown(f"- **Xu hướng tỉ lệ sinh:** {r['Xu hướng tỉ lệ sinh']}")
    st.markdown("---")

# ---------------- Export functions (safe handling)
st.markdown("### 📤 Xuất báo cáo")

col1, col2 = st.columns(2)

with col1:
    if st.button("📄 Xuất báo cáo PDF"):
        with st.spinner("Đang tạo file PDF..."):
            time.sleep(0.8)
            # Build PDF in temp file then read bytes and delete
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmpfile:
                doc = SimpleDocTemplate(tmpfile.name, pagesize=A4)
                styles = getSampleStyleSheet()
                styles.add(ParagraphStyle(name="CenterTitle", alignment=TA_CENTER, fontSize=16))
                Story = [
                    Paragraph(f"BÁO CÁO PHÂN TÍCH DÂN SỐ - {' vs '.join(sel_names)}", styles["CenterTitle"]),
                    Spacer(1, 12),
                    Paragraph(f"Giai đoạn: {start_year} - {end_year}", styles["Heading2"]),
                    Spacer(1, 12),
                ]
                # Table
                table_data = [["Quốc gia", "Tỉ lệ sinh TB (‰)", "Tỉ lệ tử TB (‰)", "Xu hướng"]] + [
                    [row["Quốc gia"], row["Tỉ lệ sinh TB (‰)"], row["Tỉ lệ tử TB (‰)"], row["Xu hướng tỉ lệ sinh"]]
                    for row in report_rows
                ]
                t = Table(table_data, hAlign="LEFT")
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
                    ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                ]))
                Story.append(t)
                Story.append(Spacer(1, 12))
                # Details textual part
                for r in report_rows:
                    Story.append(Paragraph(f"<b>{r['Quốc gia']}</b>", styles["Heading3"]))
                    Story.append(Paragraph(f"Tỉ lệ sinh trung bình: {r['Tỉ lệ sinh TB (‰)']} ‰", styles["Normal"]))
                    Story.append(Paragraph(f"Tỉ lệ tử trung bình: {r['Tỉ lệ tử TB (‰)']} ‰", styles["Normal"]))
                    Story.append(Paragraph(f"Xu hướng: {r['Xu hướng tỉ lệ sinh']}", styles["Normal"]))
                    Story.append(Spacer(1, 6))
                Story.append(Spacer(1, 12))
                Story.append(Paragraph(f"Ngày tạo: {datetime.now().strftime('%d/%m/%Y')}", styles["Normal"]))

                doc.build(Story)
                pdf_path = tmpfile.name

            try:
                with open(pdf_path, "rb") as f:
                    pdf_bytes = f.read()
            finally:
                try:
                    os.remove(pdf_path)
                except Exception:
                    pass

            st.success("✅ Đã tạo xong báo cáo PDF!")
            st.download_button(
                label="📥 Tải về báo cáo",
                data=pdf_bytes,
                file_name=f"Bao_cao_dan_so_{'_'.join(sel_names)}_{start_year}_{end_year}.pdf",
                mime="application/pdf",
            )

with col2:
    if st.button("📊 Xuất slide PowerPoint (PPTX)"):
        with st.spinner("Đang tạo slide..."):
            time.sleep(0.8)
            prs = Presentation()
            # Title slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "BÁO CÁO PHÂN TÍCH DÂN SỐ"
            slide.placeholders[1].text = f"Giai đoạn: {start_year}-{end_year}\n" + ", ".join(sel_names)

            # One slide per country
            for r in report_rows:
                s = prs.slides.add_slide(prs.slide_layouts[1])
                s.shapes.title.text = f"Tổng quan - {r['Quốc gia']}"
                s.placeholders[1].text = (
                    f"Tỉ lệ sinh TB: {r['Tỉ lệ sinh TB (‰)']}‰\n"
                    f"Tỉ lệ tử TB: {r['Tỉ lệ tử TB (‰)']}‰\n"
                    f"Xu hướng: {r['Xu hướng tỉ lệ sinh']}"
                )

            # Save to temp and provide download
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmpf:
                prs.save(tmpf.name)
                pptx_path = tmpf.name

            try:
                with open(pptx_path, "rb") as f:
                    pptx_bytes = f.read()
            finally:
                try:
                    os.remove(pptx_path)
                except Exception:
                    pass

            st.success("✅ Đã tạo file PPTX!")
            st.download_button(
                label="📥 Tải PPTX",
                data=pptx_bytes,
                file_name=f"Bao_cao_slides_{'_'.join(sel_names)}_{start_year}_{end_year}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )

# ---------------- AI Analysis tab (new)
st.markdown("---")
st.header("🧠 Phân tích chuyên sâu (AI)")

st.markdown("Chọn model, nhập API key và nhấn **Phân tích**. Hệ thống sẽ gửi tóm tắt dữ liệu (hiện tại) tới model để nhận báo cáo dài.")

col_a, col_b = st.columns([2, 3])
with col_a:
    model_choice = st.selectbox("Chọn model:", options=["openai", "gemini"])
    api_key_input = st.text_input("API key / token (ẩn):", type="password")
    analyze_button = st.button("🚀 Phân tích chuyên sâu bằng AI")

with col_b:
    st.info("Lưu ý:\n- OpenAI yêu cầu package `openai`.\n- Gemini yêu cầu `google-generativeai`.\nNếu SDK chưa cài, bạn vẫn có thể nhập key nhưng lời gọi sẽ báo lỗi.")

# Prepare prompt summary from report_rows
summary_text = ""
for r in report_rows:
    summary_text += f"{r['Quốc gia']}: Sinh TB {r['Tỉ lệ sinh TB (‰)']}‰, Tử TB {r['Tỉ lệ tử TB (‰)']}‰, Xu hướng {r['Xu hướng tỉ lệ sinh']}\n"

if analyze_button:
    if not api_key_input:
        st.error("Vui lòng nhập API key/token trước khi phân tích.")
    elif not report_rows:
        st.error("Không có dữ liệu báo cáo để gửi phân tích.")
    else:
        with st.spinner("🤖 Đang gọi model để phân tích..."):
            time.sleep(0.8)
            prompt = (
                "Bạn là chuyên gia nhân khẩu học. Dựa trên tóm tắt sau, viết một báo cáo phân tích chi tiết bằng tiếng Việt.\n\n"
                f"Tóm tắt dữ liệu:\n{summary_text}\n\n"
                "Yêu cầu:\n"
                "- Viết phần Giới thiệu, Phân tích xu hướng, Lợi ích, Rủi ro, Biện pháp ngay, Dự kiến 1-10 năm, Kết luận.\n"
                "- Trình bày rõ, có đề mục, 600-1000 từ nếu có thể.\n"
            )
            analysis = None

            # ✅ OPENAI (phiên bản >=1.0.0)
            if model_choice == "openai":
                if not OPENAI_AVAILABLE:
                    st.error("openai SDK chưa cài. Cài bằng: pip install openai")
                else:
                    try:
                        from openai import OpenAI
                        client = OpenAI(api_key=api_key_input)
                        response = client.chat.completions.create(
                            model="gpt-4o-mini",
                            messages=[
                                {"role": "system", "content": "Bạn là chuyên gia nhân khẩu học."},
                                {"role": "user", "content": prompt}
                            ],
                            max_tokens=1500,
                        )
                        analysis = response.choices[0].message.content
                    except Exception as e:
                        st.error(f"Lỗi khi gọi OpenAI: {e}")

            # ✅ GEMINI (google-generativeai mới nhất)
            else:
                if not GEMINI_AVAILABLE:
                    st.error("Gemini SDK chưa cài. Cài bằng: pip install google-generativeai")
                else:
                    try:
                        import google.generativeai as genai
                        genai.configure(api_key=api_key_input)
                        model = genai.GenerativeModel("gemini-2.5-flash")
                        response = model.generate_content(prompt)
                        analysis = response.text
                    except Exception as e:
                        st.error(f"Lỗi khi gọi Gemini: {e}")

            if analysis:
                st.success("✅ Phân tích AI hoàn tất")
                st.text_area("📜 Báo cáo phân tích AI (dài)", value=analysis, height=400)
                st.download_button("⬇️ Tải báo cáo AI (.txt)", data=analysis, file_name="Bao_cao_AI.txt", mime="text/plain")

st.caption('Gợi ý: pip install openai google-generativeai (tuỳ model bạn muốn dùng).')
